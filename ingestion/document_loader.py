"""Document loader using direct parsers: pypdf, python-pptx, pytesseract+Pillow, open()."""

from pathlib import Path
from langchain_core.documents import Document
from pypdf import PdfReader
from pptx import Presentation
from PIL import Image
import pytesseract

from utils.helpers import utcnow_iso, clean_text, file_extension
from utils.logger import get_logger

log = get_logger("document_loader")

# Supported file types → parser map
_IMAGE_EXTS = {"png", "jpg", "jpeg", "tiff", "bmp", "gif", "webp"}
_TEXT_EXTS = {"txt", "md", "markdown", "csv", "json", "rst"}


def load_document(file_path: str, filename: str) -> list[Document]:
    """
    Parse a single file and return a list of LangChain Documents.

    Supported formats:
        - PDF  → pypdf
        - PPTX → python-pptx
        - Images (png/jpg/…) → pytesseract OCR
        - Text / Markdown → plain open()
    """
    ext = file_extension(filename)
    log.info(f"Loading {filename} (type={ext})")

    try:
        if ext == "pdf":
            return _load_pdf(file_path, filename)
        elif ext == "pptx":
            return _load_pptx(file_path, filename)
        elif ext in _IMAGE_EXTS:
            return _load_image(file_path, filename)
        elif ext in _TEXT_EXTS:
            return _load_text(file_path, filename)
        else:
            log.warning(f"Unsupported file type: {ext} — attempting plain text read")
            return _load_text(file_path, filename)
    except Exception as e:
        log.error(f"Failed to load {filename}: {e}")
        return []


def _make_metadata(filename: str, file_type: str, page_number: int = 0) -> dict:
    """Build the base metadata dict for a document."""
    return {
        "source_file": filename,
        "file_type": file_type,
        "page_number": page_number,
        "element_type": "NarrativeText",
        "subject": "",
        "chapter": "",
        "difficulty": "intermediate",
        "ingestion_timestamp": utcnow_iso(),
    }


# ── PDF ──────────────────────────────────────────────────────────────
def _load_pdf(file_path: str, filename: str) -> list[Document]:
    reader = PdfReader(file_path)
    docs: list[Document] = []

    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        text = clean_text(text)

        if not text.strip():
            # Attempt OCR on pages with no extracted text (scanned pages)
            try:
                from pypdf import PageObject
                # Try to extract images for OCR
                for image_obj in page.images:
                    img = Image.open(image_obj.data)
                    ocr_text = pytesseract.image_to_string(img)
                    if ocr_text.strip():
                        text += "\n" + clean_text(ocr_text)
            except Exception:
                pass

        if text.strip():
            meta = _make_metadata(filename, "pdf", page_number=page_num)
            docs.append(Document(page_content=text, metadata=meta))

    log.info(f"PDF {filename}: extracted {len(docs)} pages")
    return docs


# ── PPTX ─────────────────────────────────────────────────────────────
def _load_pptx(file_path: str, filename: str) -> list[Document]:
    prs = Presentation(file_path)
    docs: list[Document] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        texts.append(para_text)

            # Also try to extract text from tables
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        texts.append(" | ".join(row_texts))

        combined = clean_text("\n".join(texts))
        if combined.strip():
            meta = _make_metadata(filename, "pptx", page_number=slide_num)
            meta["element_type"] = "SlideContent"
            docs.append(Document(page_content=combined, metadata=meta))

    log.info(f"PPTX {filename}: extracted {len(docs)} slides")
    return docs


# ── Image (OCR) ──────────────────────────────────────────────────────
def _load_image(file_path: str, filename: str) -> list[Document]:
    img = Image.open(file_path)
    text = pytesseract.image_to_string(img)
    text = clean_text(text)

    if not text.strip():
        log.warning(f"Image {filename}: OCR produced no text")
        return []

    meta = _make_metadata(filename, "image")
    meta["element_type"] = "OCRText"
    log.info(f"Image {filename}: extracted {len(text)} chars via OCR")
    return [Document(page_content=text, metadata=meta)]


# ── Plain text / Markdown ────────────────────────────────────────────
def _load_text(file_path: str, filename: str) -> list[Document]:
    path = Path(file_path)
    text = path.read_text(encoding="utf-8", errors="replace")
    text = clean_text(text)

    if not text.strip():
        return []

    ext = file_extension(filename)
    file_type = "markdown" if ext in ("md", "markdown") else "text"
    meta = _make_metadata(filename, file_type)
    meta["element_type"] = "PlainText"
    log.info(f"Text {filename}: {len(text)} chars")
    return [Document(page_content=text, metadata=meta)]
